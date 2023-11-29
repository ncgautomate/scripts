from pptx import Presentation
from pptx.util import Pt
import os

def read_markdown_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def add_slide(presentation, title, content):
    slide_layout = presentation.slide_layouts[5]  # Choose a slide layout
    slide = presentation.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    content_placeholder.text = content

# Path to your Markdown file
#md_file_path = os.path.join('E:', '\\Github\\orgs\\ncgcloudhub\\scripts\\word-docs-to-pptmarkdown-to-ppt-content.md')  # Update this to the path of your .md file
#md_file_path = r'E:\Github\orgs\ncgcloudhub\scripts\word-docs-to-ppt\markdown-to-ppt-content.md'
md_file_path = os.path.join('E:\\', 'Github', 'orgs', 'ncgcloudhub', 'scripts', 'word-docs-to-ppt', 'markdown-to-ppt-content.md')


# Read the contents of the Markdown file
markdown_content = read_markdown_file(md_file_path)

# Process the Markdown content to extract slide information
# This part depends on how the slides are formatted in your Markdown
# Assuming each slide is separated by a line "---"

slides = markdown_content.split('---')
prs = Presentation()

for slide_content in slides:
    # You need to parse the slide_content to extract the title and the rest of the content
    # This is a placeholder for the actual parsing logic
    title, content = slide_content.strip().split('\n', 1)
    add_slide(prs, title.strip('# '), content)

# Save the presentation
presentation_file_path = os.path.join('E:', 'path_to_your_presentation.pptx')  # Update this to your desired path
prs.save(presentation_file_path)
