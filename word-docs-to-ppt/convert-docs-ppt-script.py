from pptx import Presentation
from docx import Document

def create_presentation_from_docx(docx_path, pptx_path):
    # Load the Word document
    doc = Document(docx_path)
    
    # Create a PowerPoint presentation object
    prs = Presentation()
    
    # Iterate through the document to find titles that start with "Slide"
    for paragraph in doc.paragraphs:
        if paragraph.text.startswith("Slide"):
            # Add a slide
            slide_layout = prs.slide_layouts[1]  # 1 is for 'Title and Content'
            slide = prs.slides.add_slide(slide_layout)
            
            # Set the title and content for the slide
            title, content = paragraph.text.split(": ", 1)
            slide.shapes.title.text = content  # Assuming the slide title follows 'Slide X: Title'
            slide.placeholders[1].text = "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text)
    
    # Save the presentation
    prs.save(pptx_path)

# Example usage
docx_file_path = 'path_to_your_document.docx'
pptx_file_path = 'desired_path_for_presentation.pptx'
create_presentation_from_docx(docx_file_path, pptx_file_path)
